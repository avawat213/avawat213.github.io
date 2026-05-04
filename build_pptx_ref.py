#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

LIGHT_PINK = RGBColor(255, 217, 236)
MEDIUM_PINK = RGBColor(255, 194, 226)
DARK_PINK = RGBColor(94, 17, 65)
GOLD = RGBColor(255, 215, 0)

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide1.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = LIGHT_PINK

title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Title"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = DARK_PINK
p.alignment = PP_ALIGN.CENTER

sub_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "Subtitle"
p.font.size = Pt(28)
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.CENTER

# Slide 2: Bullet Slide
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide2.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = LIGHT_PINK

head_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
tf = head_box.text_frame
p = tf.paragraphs[0]
p.text = "Heading 1"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = DARK_PINK

cont_box = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
tf = cont_box.text_frame
tf.word_wrap = True

for i, text in enumerate(["Bullet point one", "Bullet point two", "Bullet point three"]):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_PINK
    p.space_before = Pt(6)
    p.space_after = Pt(6)

prs.save('slide-style.pptx')
print("Reference PowerPoint created successfully!")
